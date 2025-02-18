from flask import Flask, render_template, request, send_file
import requests
from pptx import Presentation
from pptx.util import Inches
import io
from PIL import Image as PILImage
import re
import os
import google.generativeai as genai
import random

app = Flask(__name__)

# Load environment variables from .env file

# Set your OpenAI API key
genai.configure(api_key="AIzaSyA3Xz0Q-7JsDx0qrPuUhnIJJwifttSwtIE")
# Set your Pexels API key
# PEXELS_API_KEY = "XZmXIRIMHkUT9ZHKN9IUpxQ8INRxd9kjsKAKWuIIcQ3knYaI9ha1h5iE"

@app.route('/', methods=['GET'])
def index():
    return render_template('index.html')

@app.route('/generate-ppt', methods=['POST'])
def generate_ppt():
    title = request.form['title']
    presenter = request.form['presenter']
    num_slides = int(request.form['num_slides'])
    prompt = request.form['prompt']
    include_references = 'include_references' in request.form
    include_images = 'include_images' in request.form
    template_choice = request.form['template_choice']
    template_path = f'presentations/{template_choice}'

    slides_content, conclusion_content, references_content = generate_content(
        prompt, num_slides, include_references
    )
    ppt_file = create_ppt(slides_content, conclusion_content, references_content, title, presenter, template_path, include_images, num_slides)
    return send_file(ppt_file, as_attachment=True, download_name='presentation.pptx')

def generate_content(prompt, num_slides, include_references):


    # Initialize the Gemini model for main content
    model = genai.GenerativeModel(model_name="gemini-2.0-flash")
    # Generate main slide content
    response_main = model.generate_content(f"{prompt}\n\nGenerate content for {num_slides} slides, formatted as separate paragraphs.")
    
    main_content = response_main.text.strip()
    slides_content = main_content.split("\n\n")[:num_slides]  # Split content into slides

    # Use a more advanced model for summarization and references
    model_pro = genai.GenerativeModel(model_name="gemini-pro")

    # Generate conclusion
    response_conclusion = model_pro.generate_content(f"Summarize the key points covered in the following content:\n{main_content}")
    conclusion_content = response_conclusion.text.strip()

    # Generate references (if needed)
    references_content = ""
    if include_references:
        response_references = model_pro.generate_content("Provide a list of references based on the discussed topics.")
        references_content = response_references.text.strip()

    return slides_content, conclusion_content, references_content



# def fetch_image(query):
#     headers = {'Authorization': PEXELS_API_KEY}
#     params = {'query': query, 'per_page': 1}
#     response = requests.get('https://api.pexels.com/v1/search', headers=headers, params=params)
#     if response.status_code == 200:
#         results = response.json()
#         if results['photos']:
#             image_url = results['photos'][0]['src']['original']
#             image_response = requests.get(image_url)
#             if image_response.status_code == 200:
#                 image = PILImage.open(io.BytesIO(image_response.content))
#                 return image
#     return None

# def fetch_image(query):
#     if not PEXELS_API_KEY:
#         print("Error: Pexels API key is missing.")
#         return None

#     headers = {'Authorization': PEXELS_API_KEY}
#     params = {'query': request.form['title'], 'per_page': 1}

#     response = requests.get('https://api.pexels.com/v1/search', headers=headers, params=params)

#     if response.status_code == 200:
#         results = response.json()
#         if results['photos']:
#             image_url = results['photos'][0]['src']['original']
#             image_response = requests.get(image_url)

#             if image_response.status_code == 200:
#                 image = PILImage.open(io.BytesIO(image_response.content))
#                 print("✅ Image fetched successfully!")
#                 image.show()  # Display the image
#                 return image

#     print("⚠️ No image found or API request failed.")
#     return None

def fetch_image(query):
    if not PEXELS_API_KEY:
        print("Error: Pexels API key is missing.")
        return None

    headers = {'Authorization': PEXELS_API_KEY}
    params = {'query': request.form['title'], 'per_page': 5}  # Fetch multiple images

    response = requests.get('https://api.pexels.com/v1/search', headers=headers, params=params)

    if response.status_code == 200:
        results = response.json()
        if results['photos']:
            random_image = random.choice(results['photos'])  # Pick a random image
            image_url = random_image['src']['original']
            image_response = requests.get(image_url)

            if image_response.status_code == 200:
                image = PILImage.open(io.BytesIO(image_response.content))
                print(f"✅ Random image fetched successfully: {image_url}")
                image.show()  # Display the image
                return image

    print("⚠️ No image found or API request failed.")
    return None

def add_image_to_slide(slide, image, prs):
    image_stream = io.BytesIO()
    image.save(image_stream, format='PNG')
    image_stream.seek(0)
    slide.shapes.add_picture(image_stream, prs.slide_width - Inches(4), Inches(1), width=Inches(4), height=Inches(3))

def create_ppt(slides_content, conclusion_content, references_content, presentation_title, presenter_name, template_path, include_images, num_slides):
    prs = Presentation(template_path)

    # Remove all existing slides from the presentation
    while len(prs.slides) > 0:
        xml_slides = prs.slides._sldIdLst
        prs.part.drop_rel(xml_slides[0].rId)
        del xml_slides[0]

    # Create the title slide
    title_slide_layout = prs.slide_layouts[0]  # Assuming the first layout is for the title
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = presentation_title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Presented by " + presenter_name

    # Create content slides using the specified content layout
    content_slide_layout = prs.slide_layouts[2]  # TITLE_AND_BODY layout
    for content in slides_content[:num_slides]:
        slide = prs.slides.add_slide(content_slide_layout)
        parts = content.split(':', 1)
        title = re.sub(r'\*\*|__|```', '', parts[0]).strip()  # Clean markdown syntax
        body = parts[1].strip() if len(parts) > 1 else ""

        if len(title.split()) > 5:
            body = title + " " + body  # Long title, move to body
            title = "Overview"

        slide.shapes.title.text = title
        slide.placeholders[1].text = body  # Assumed body placeholder exists in this layout

        if include_images and len(title.split()) <= 5:
            image = fetch_image(title)
            if image:
                add_image_to_slide(slide, image, prs)

    # Optional conclusion slide
    if conclusion_content:
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = "Conclusion"
        slide.placeholders[1].text = conclusion_content

    # Optional references slide
    if references_content:
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = "References"
        slide.placeholders[1].text = references_content

    file_path = 'generated_presentation_using_template.pptx'
    prs.save(file_path)
    return file_path




if __name__ == '__main__':
    app.run(debug=True)
