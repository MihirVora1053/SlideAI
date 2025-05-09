import streamlit as st
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
from pydub import AudioSegment
import speech_recognition as sr
import os
import io
import json
from flask import Flask, request, send_file
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import scraping1  # Importing refined transcript generation

# Load the pre-trained tokenizer and model
@st.cache_resource
def load_model():
    tokenizer = AutoTokenizer.from_pretrained("facebook/bart-large-cnn")
    model = AutoModelForSeq2SeqLM.from_pretrained("facebook/bart-large-cnn")
    return tokenizer, model

tokenizer, model = load_model()
AudioSegment.converter = "/path/to/ffmpeg"


def process_audio(file_path):
    try:
        temp_file_path = "temp_uploaded_audio.wav"
        with open(temp_file_path, "wb") as f:
            f.write(file_path.read())

        audio = AudioSegment.from_file(temp_file_path)
        mono_audio = audio.set_channels(1)
        converted_audio_path = "temp_audio.wav"
        mono_audio.export(converted_audio_path, format="wav")

        recognizer = sr.Recognizer()
        transcribed_text = ""

        # Split into 30-second chunks
        chunk_length_ms = 30 * 1000  # 30 seconds
        chunks = [mono_audio[i:i + chunk_length_ms] for i in range(0, len(mono_audio), chunk_length_ms)]

        for i, chunk in enumerate(chunks):
            chunk_path = f"chunk_{i}.wav"
            chunk.export(chunk_path, format="wav")

            with sr.AudioFile(chunk_path) as source:
                audio_data = recognizer.record(source)
                try:
                    text = recognizer.recognize_google(audio_data)
                    transcribed_text += text + " "
                except sr.UnknownValueError:
                    continue  # Skip unrecognized chunks
                except sr.RequestError as e:
                    return f"API error during chunk {i}: {e}"

            os.remove(chunk_path)

        return transcribed_text.strip()
    except Exception as e:
        return f"Error processing audio: {e}"


def generate_ppt_and_notes(transcript, selected_template):    
    template_path = f"presentations/{selected_template}"
    filepath = "slides.json"
    notes_file = scraping1.generate_notes(transcript, summary_length)

    with open(filepath, "r") as file:
        slides_data = json.load(file)
    ppt_file = create_ppt(slides_data, template_path)
    
    return ppt_file, notes_file

def create_ppt(slide_data, template_path):
    prs = Presentation(template_path)
    
    while len(prs.slides) > 0:
        xml_slides = prs.slides._sldIdLst
        prs.part.drop_rel(xml_slides[0].rId)
        del xml_slides[0]

    title_color = RGBColor(255, 69, 0)
    text_color = RGBColor(50, 50, 50)
    background_color = RGBColor(240, 240, 240)

    max_chars_per_slide = 550  # Adjust content limit per slide
    
    for slide_content in slide_data:
        title = slide_content["title"]
        points = slide_content["points"]
        
        while points:
            slide_layout = prs.slide_layouts[2]  # TITLE_AND_BODY layout
            slide = prs.slides.add_slide(slide_layout)

            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = background_color

            slide_title = slide.shapes.title
            slide_title.text = title
            slide_title.text_frame.paragraphs[0].font.size = Pt(32)
            slide_title.text_frame.paragraphs[0].font.bold = True
            slide_title.text_frame.paragraphs[0].font.color.rgb = title_color

            if len(slide.placeholders) > 1:
                content = slide.placeholders[1].text_frame
                remaining_points = []

                char_count = 0
                for point in points:
                    char_count += len(point)

                    if char_count > max_chars_per_slide:
                        remaining_points.append(point)
                    else:
                        p = content.add_paragraph()
                        p.text = point
                        p.font.size = Pt(16)
                        p.font.color.rgb = text_color
                        p.font.name = "Calibri"

                points = remaining_points  # Move extra points to next slide

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

# Function to list available PPT templates
def get_ppt_templates():
    template_folder = "presentations"
    return [f for f in os.listdir(template_folder) if f.endswith(".pptx")]

st.title("Automated Lecture Processing")
st.write("Upload an audio file to generate lecture notes and presentation.")

uploaded_file = st.file_uploader("Upload Audio File (WAV format)", type=["wav"])

# Dropdown for summary length selection
summary_length = st.selectbox(
    "Select Summary Length",
    ["Short", "Medium", "Detailed"]
)

# PPT template selection
ppt_templates = get_ppt_templates()
selected_template = st.selectbox("Choose PPT Template:", ppt_templates if ppt_templates else ["No templates found"])

if "ppt_file" not in st.session_state:
    st.session_state.ppt_file = None
if "notes_file" not in st.session_state:
    st.session_state.notes_file = None

# Button to start the process
if st.button("Generate PPT"):
    if uploaded_file:
        with st.spinner("Processing audio..."):
            transcribed_text = process_audio(uploaded_file)
            if transcribed_text and not transcribed_text.startswith("Error"):
                ppt_file, notes_file = generate_ppt_and_notes(transcribed_text, selected_template)

                # Store files in session state
                st.session_state.ppt_file = ppt_file
                st.session_state.notes_file = notes_file
                st.success("Lecture notes and presentation generated successfully!")
            else:
                st.error(transcribed_text)
    else:
        st.warning("Please upload an audio file before generating the PPT.")

# Display download buttons only if files exist
if st.session_state.ppt_file:
    st.download_button("Download Presentation", st.session_state.ppt_file, file_name="lecture_notes.pptx")

if st.session_state.notes_file:
    st.download_button("Download Notes", st.session_state.notes_file, file_name="lecture_notes.docx")