from flask import Flask, render_template, request, send_file
import json
import re
import nltk
from nltk.tokenize import sent_tokenize
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import io
import scraping

nltk.download("punkt")  # Ensure sentence tokenizer is available

app = Flask(__name__)


def transcript_to_json(transcript):
    slides = []
    sections = re.split(
        r"Topic: (.+)", transcript
    )  # Removed \n to include the first topic

    for i in range(1, len(sections), 2):
        title = sections[i].strip()
        content = sections[i + 1].strip()
        sentences = sent_tokenize(content)
        slides.append({"title": title, "points": sentences})

    return slides


# def create_ppt(slide_data, template_path):
#     prs = Presentation(template_path)

#     # Remove all existing slides from the presentation
#     while len(prs.slides) > 0:
#         xml_slides = prs.slides._sldIdLst
#         prs.part.drop_rel(xml_slides[0].rId)
#         del xml_slides[0]

#     # Create the title slide
#     title_slide_layout = prs.slide_layouts[
#         0
#     ]  # Assuming the first layout is for the title
#     slide = prs.slides.add_slide(title_slide_layout)
#     # slide.shapes.title.text = presentation_title
#     # if len(slide.placeholders) > 1:
#         # slide.placeholders[1].text = "Presented by " + presenter_name

#     # Create content slides using the specified content layout
#     content_slide_layout = prs.slide_layouts[1]  # TITLE_AND_BODY layout
#     # for content in slide_data[:num_slides]:
#     #     slide = prs.slides.add_slide(content_slide_layout)
#     #     parts = content.split(":", 1)
#     #     title = re.sub(r"\*\*|__|```", "", parts[0]).strip()  # Clean markdown syntax
#     #     body = parts[1].strip() if len(parts) > 1 else ""

#     #     if len(title.split()) > 5:
#     #         body = title + " " + body  # Long title, move to body
#     #         title = "Overview"

#     #     slide.shapes.title.text = title
#     #     slide.placeholders[1].text = (
#     #         body  # Assumed body placeholder exists in this layout
#     #     )

#     #     if include_images and len(title.split()) <= 5:
#     #         image = fetch_image(title)
#     #         if image:
#     #             add_image_to_slide(slide, image, prs)

#     # Optional conclusion slide
#     # if conclusion_content:
#     #     slide = prs.slides.add_slide(content_slide_layout)
#     #     slide.shapes.title.text = "Conclusion"
#     #     slide.placeholders[1].text = conclusion_content

#     # # Optional references slide
#     # if references_content:
#     #     slide = prs.slides.add_slide(content_slide_layout)
#     #     slide.shapes.title.text = "References"
#     #     slide.placeholders[1].text = references_content

#     title_color = RGBColor(255, 69, 0)
#     text_color = RGBColor(50, 50, 50)
#     background_color = RGBColor(240, 240, 240)

#     for slide_content in slide_data:
#         slide_layout = prs.slide_layouts[1]
#         slide = prs.slides.add_slide(slide_layout)

#         background = slide.background
#         fill = background.fill
#         fill.solid()
#         fill.fore_color.rgb = background_color

#         title = slide.shapes.title
#         title.text = slide_content["title"]
#         title.text_frame.paragraphs[0].font.size = Pt(32)
#         title.text_frame.paragraphs[0].font.bold = True
#         title.text_frame.paragraphs[0].font.color.rgb = title_color


#         if len(slide.placeholders) > 1:
#             content = slide.placeholders[1].text_frame
#         # else:
#         #     print("⚠️ Placeholder[1] missing, adding text box manually.")
#         #     content = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5)).text_frame

#         for point in slide_content["points"]:
#             p = content.add_paragraph()
#             p.text = point
#             p.font.size = Pt(20)
#             p.font.color.rgb = text_color


#     ppt_stream = io.BytesIO()
#     prs.save(ppt_stream)
#     ppt_stream.seek(0)
#     return ppt_stream


# @app.route("/", methods=["GET"])
# def index():
#     return render_template("index.html")


# @app.route("/generate-ppt", methods=["POST"])
# def generate_ppt():
#     template_choice = request.form["template_choice"]
#     template_path = f"presentations/{template_choice}"
#     transcript = request.form["prompt"]
#     slides_data = transcript_to_json(transcript)

#     with open("slides.json", "w") as f:
#         json.dump(slides_data, f, indent=4)

#     ppt_file = create_ppt(slides_data, template_path)
#     return send_file(
#         ppt_file,
#         as_attachment=True,
#         download_name="lecture_notes.pptx",
#         mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
#     )


def create_ppt(slide_data, template_path):
    prs = Presentation(template_path)

    # Remove existing slides
    while len(prs.slides) > 0:
        xml_slides = prs.slides._sldIdLst
        prs.part.drop_rel(xml_slides[0].rId)
        del xml_slides[0]

    title_color = RGBColor(255, 69, 0)
    text_color = RGBColor(50, 50, 50)
    background_color = RGBColor(240, 240, 240)

    max_chars_per_slide = 900  # Adjust as needed for fitting content

    for slide_content in slide_data:
        title = slide_content["title"]
        points = slide_content["points"]

        while points:
            slide_layout = prs.slide_layouts[2]  # TITLE_AND_BODY layout
            slide = prs.slides.add_slide(slide_layout)

            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = background_color

            slide_title = slide.shapes.title
            slide_title.text = title
            slide_title.text_frame.paragraphs[0].font.size = Pt(32)
            slide_title.text_frame.paragraphs[0].font.bold = True
            slide_title.text_frame.paragraphs[0].font.color.rgb = title_color

            if len(slide.placeholders) > 1:
                content = slide.placeholders[1].text_frame
                remaining_points = []

                char_count = 0
                for point in points:
                    char_count += len(point)

                    if (
                        char_count > max_chars_per_slide
                    ):  # If exceeds space, move to next slide
                        remaining_points.append(point)
                    else:
                        p = content.add_paragraph()
                        p.text = point
                        p.font.size = Pt(14)
                        p.font.color.rgb = text_color
                        p.font.name = "Calibri"

                points = remaining_points  # Move extra points to the next slide
                # title = "Continued: " + title  # Update title for continuity

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream


app = Flask(__name__)


@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")


@app.route("/generate-ppt", methods=["POST"])
def generate_ppt():
    template_choice = request.form["template_choice"]
    template_path = f"presentations/{template_choice}"
    transcript = request.form["prompt"]
    
    # transcript = request.form["prompt"]
    slides_data = transcript_to_json(transcript)

    with open("slides.json", "w") as f:
        json.dump(slides_data, f, indent=4)

    ppt_file = create_ppt(slides_data, template_path)
    return send_file(
        ppt_file,
        as_attachment=True,
        download_name="lecture_notes.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

# print(scraping.generate_notes.refined_key_points)
if __name__ == "__main__":
    app.run(debug=True)
