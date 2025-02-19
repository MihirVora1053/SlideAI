from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import json
import re
import nltk
from nltk.tokenize import sent_tokenize

nltk.download("punkt_tab")  # Download sentence tokenizer if not already downloaded


def transcript_to_json(transcript):
    slides = []

    # Split transcript into sections based on headings (assuming they start with "Topic: ")
    sections = re.split(r"\nTopic: (.+)", transcript)

    # Iterate over sections and extract structured data
    for i in range(1, len(sections), 2):  # Skip first empty split result
        title = sections[i].strip()  # Extract title
        content = sections[i + 1].strip()  # Extract content

        # Tokenize sentences and pick key points (first few sentences)
        sentences = sent_tokenize(content)
        key_points = sentences[
            : min(5, len(sentences))
        ]  # Limit key points to 5 per slide

        # Append structured data
        slides.append({"title": title, "points": key_points})

    return slides  # Returning list instead of JSON string


# Example transcript
transcript = """
Topic: Income Tax
The income tax in India is a direct tax levied using a progressive slab rate system, where the tax rate increases as the taxpayer's income rises.  
The income tax in India is a direct tax levied using a progressive slab rate system, where the tax rate increases as the taxpayer's income rises.  
The income tax in India is a direct tax levied using a progressive slab rate system, where the tax rate increases as the taxpayer's income rises.  
The Income Tax Act of 1961 provides two tax regimes: the old regime (with deductions and exemptions) and the new regime (with lower tax rates but fewer deductions). 
This article will explain the income tax calculation under both regimes, comparing tax slabs and highlighting key differences.  
Significant changes were introduced in the new regime, applicable for the financial year 2024-25, including revised slabs, an increased standard deduction of Rs 75,000 for salaried employees, and a higher NPS contribution limit (increased to 14% from 10%).  
Taxpayers can now choose between the two regimes to maximize their tax benefits.
The income tax in India is a direct tax levied using a progressive slab rate system, where the tax rate increases as the taxpayer's income rises.  
The income tax in India is a direct tax levied using a progressive slab rate system, where the tax rate increases as the taxpayer's income rises.  
The income tax in India is a direct tax levied using a progressive slab rate system, where the tax rate increases as the taxpayer's income rises.  
The Income Tax Act of 1961 provides two tax regimes: the old regime (with deductions and exemptions) and the new regime (with lower tax rates but fewer deductions). 
This article will explain the income tax calculation under both regimes, comparing tax slabs and highlighting key differences.  
Significant changes were introduced in the new regime, applicable for the financial year 2024-25, including revised slabs, an increased standard deduction of Rs 75,000 for salaried employees, and a higher NPS contribution limit (increased to 14% from 10%).  
Taxpayers can now choose between the two regimes to maximize their tax benefits.

Topic: Introduction to AI
Artificial Intelligence (AI) is the simulation of human intelligence in machines. 
It enables machines to learn from experience and perform tasks that typically require human cognition. 
There are different types of AI, including Narrow AI and General AI. 
AI applications are found in various domains, such as healthcare, finance, and robotics.

Topic: Machine Learning
Machine learning is a subset of AI that allows computers to learn from data without being explicitly programmed. 
It consists of three main types: supervised learning, unsupervised learning, and reinforcement learning. 
Each type of machine learning has its own unique characteristics and applications.
"""

# Convert transcript to JSON format
slides_data = transcript_to_json(transcript)

# Save JSON to file (optional)
with open("slides.json", "w") as f:
    json.dump(slides_data, f, indent=4)

print(json.dumps(slides_data, indent=4))


def create_ppt(slide_data, filename="lecture_notes.pptx"):
    # Create a presentation object
    prs = Presentation()

    # Define custom colors
    title_color = RGBColor(255, 69, 0)  # Orange-Red
    text_color = RGBColor(50, 50, 50)  # Dark Gray
    background_color = RGBColor(240, 240, 240)  # Light Gray

    for slide_content in slide_data:
        # Add a new slide with a title & content layout
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # Set slide background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color

        # Set title with custom font size and color
        title = slide.shapes.title
        title.text = slide_content["title"]
        title.text_frame.text = slide_content["title"]
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = title_color

        # Set bullet points with custom font size and color
        content = slide.placeholders[1].text_frame
        for point in slide_content["points"]:
            p = content.add_paragraph()
            p.text = point
            p.font.size = Pt(20)
            p.font.color.rgb = text_color

    # Save the presentation
    prs.save(filename)
    print(f"PPT file saved as {filename}")


# Generate PPT from structured slides data
create_ppt(slides_data)
