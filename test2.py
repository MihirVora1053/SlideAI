from flask import Flask, render_template, request, send_file
import json
import re
import nltk
from nltk.tokenize import sent_tokenize
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import io
import scraping1  # Importing refined transcript generation

nltk.download("punkt")  # Ensure sentence tokenizer is available

app = Flask(__name__)

# # Hardcoded transcript (Replace this later with actual audio processing)
# original_transcript = """The topic of artificial intelligence has evolved significantly over the past few decades.
# AI is now widely used in various industries such as healthcare, finance, and automation.
# Machine learning and deep learning have revolutionized data analysis.
# However, ethical concerns regarding AI biases and job automation remain critical discussions.
# The future of AI holds immense potential, with advancements in generative AI and human-like interactions."""

refined_transcript = """
Artificial Intelligence (AI) refers to the simulation of human intelligence in machines, enabling them to perform tasks that typically require cognitive abilities such as learning, reasoning, and decision-making.
Artificial Intelligence (AI) refers to the simulation of human intelligence in machines, enabling them to perform tasks that typically require cognitive abilities such as learning, reasoning, and decision-making.
Artificial Intelligence (AI) refers to the simulation of human intelligence in machines, enabling them to perform tasks that typically require cognitive abilities such as learning, reasoning, and decision-making.
AI is broadly categorized into two types: Narrow AI, which is designed to perform specific tasks like image recognition or language translation, and General AI, which can mimic human intelligence across various domains.
This article will explain the key concepts of AI, its applications across industries, and its impact on modern society.
Recent advancements in AI have led to breakthroughs in natural language processing, deep learning, and automation, transforming fields such as healthcare, finance, and robotics.
AI-powered tools like chatbots, virtual assistants, and predictive analytics are now widely used, improving efficiency and decision-making processes.
Artificial Intelligence (AI) refers to the simulation of human intelligence in machines, enabling them to perform tasks that typically require cognitive abilities such as learning, reasoning, and decision-making.
Artificial Intelligence (AI) refers to the simulation of human intelligence in machines, enabling them to perform tasks that typically require cognitive abilities such as learning, reasoning, and decision-making.
Artificial Intelligence (AI) refers to the simulation of human intelligence in machines, enabling them to perform tasks that typically require cognitive abilities such as learning, reasoning, and decision-making.
AI is broadly categorized into two types: Narrow AI, which is designed to perform specific tasks like image recognition or language translation, and General AI, which can mimic human intelligence across various domains.
This article will explain the key concepts of AI, its applications across industries, and its impact on modern society.
Recent advancements in AI have led to breakthroughs in natural language processing, deep learning, and automation, transforming fields such as healthcare, finance, and robotics.
AI-powered tools like chatbots, virtual assistants, and predictive analytics are now widely used, improving efficiency and decision-making processes."""



def create_ppt(slide_data, template_path):
    """Generates a PowerPoint presentation using structured slide data."""
    prs = Presentation(template_path)

    # Remove default slides
    while len(prs.slides) > 0:
        xml_slides = prs.slides._sldIdLst
        prs.part.drop_rel(xml_slides[0].rId)
        del xml_slides[0]

    title_color = RGBColor(255, 69, 0)
    text_color = RGBColor(50, 50, 50)
    background_color = RGBColor(240, 240, 240)

    max_chars_per_slide = 900  # Adjust content limit per slide

    for slide_content in slide_data:
        title = slide_content["title"]
        points = slide_content["points"]

        while points:
            slide_layout = prs.slide_layouts[2]  # TITLE_AND_BODY layout
            slide = prs.slides.add_slide(slide_layout)

            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = background_color

            slide_title = slide.shapes.title
            slide_title.text = title
            slide_title.text_frame.paragraphs[0].font.size = Pt(32)
            slide_title.text_frame.paragraphs[0].font.bold = True
            slide_title.text_frame.paragraphs[0].font.color.rgb = title_color

            if len(slide.placeholders) > 1:
                content = slide.placeholders[1].text_frame
                remaining_points = []

                char_count = 0
                for point in points:
                    char_count += len(point)

                    if char_count > max_chars_per_slide:
                        remaining_points.append(point)
                    else:
                        p = content.add_paragraph()
                        p.text = point
                        p.font.size = Pt(14)
                        p.font.color.rgb = text_color
                        p.font.name = "Calibri"

                points = remaining_points  # Move extra points to next slide

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream


@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")


@app.route("/generate-ppt", methods=["POST"])
def generate_ppt():
    """Generates both PPT and DOCX files for download."""

    scraping1.generate_notes(
        refined_transcript, "Artificial Intelligence"
    )  # Generates DOCX notes
    
    template_choice = request.form["template_choice"]
    template_path = f"presentations/{template_choice}"

    filepath = "slides.json"
    try:
        with open(filepath, "r") as file:
            data = json.load(file)
            # print(data)
    except FileNotFoundError:
        print(f"Error: File '{filepath}' not found.")
        return None
    except json.JSONDecodeError:
        print(f"Error: Invalid JSON format in '{filepath}'.")
        return None
    except Exception as e:
        print(f"An unexpected error occurred: {e}")
        return None
    
    slides_data = data

    ppt_file = create_ppt(slides_data, template_path)

    print(f"✅ PPT saved successfully !!!")


    return send_file(
        ppt_file,
        as_attachment=True,
        download_name="lecture_notes.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


if __name__ == "__main__":
    app.run(debug=True)
